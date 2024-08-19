import PyPDF2
from os import listdir
from os.path import isfile, join,isdir

import torch
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_qdrant import Qdrant
import sys
from langchain_text_splitters import TokenTextSplitter
from pptx import Presentation
from qdrant_client import QdrantClient
from qdrant_client.models import Distance, VectorParams
import docx
import os

def get_files(dir):
    file_list = []
    for dir, _, filenames in os.walk(dir):
        for f in filenames:
            file_list.append(os.path.join(dir, f))
    # for f in listdir(dir):
    #     if isfile(join(dir,f)):
    #         file_list.append(join(dir,f))
    #     elif isdir(join(dir,f)):
    #         file_list= file_list + get_files(join(dir,f))
    return file_list

def getTextFromWord(filename):
    doc = docx.Document(filename)
    fullText = []
    for para in doc.paragraphs:
        fullText.append(para.text)
    return '\n'.join(fullText)

def getTextFromPPTX(filename):
    prs = Presentation(filename)
    fullText = []
    for slide in prs.slides:
        for shape in slide.shapes:
            fullText.append(shape.text)
    return '\n'.join(fullText)

def main_indexing(mypath):
    #model_name = "amberoad/bert-multilingual-passage-reranking-msmarco"
    model_name = "sentence-transformers/msmarco-bert-base-dot-v5"
    if torch.cuda.is_available():
        model_kwargs = {'device': 'cuda'}
    elif torch.backends.mps.is_available():
        model_kwargs = {'device': 'mps'}
    else:
        model_kwargs = {'device': 'cpu'}
    encode_kwargs = {'normalize_embeddings': True}
    hf = HuggingFaceEmbeddings(
        model_name=model_name,
        model_kwargs=model_kwargs,
        encode_kwargs=encode_kwargs
    )
    client = QdrantClient(path="qdrant/")
    collection_name = "MyCollection"
    if client.collection_exists(collection_name):
        client.delete_collection(collection_name)

    client.create_collection(collection_name,vectors_config=VectorParams(size=768, distance=Distance.DOT))
    qdrant = Qdrant(client, collection_name, hf)
    print("Indexing...")
    onlyfiles = get_files(mypath)
    file_content = ""
    for file in onlyfiles:
        file_content = ""
        if file.find("~") > 0:  # added by pdchristian to catch files with "~" in file name
            file_content = "Empty due to ~ in file name."  # added by pdchristian to catch files with "~" in file name
            print("Document title with ~: " + file)
        elif file.endswith(".pdf"):
            try:
                print("indexing "+file)
                reader = PyPDF2.PdfReader(file)
                for i in range(0,len(reader.pages)):
                    file_content = file_content + " "+reader.pages[i].extract_text()
            except Exception as exc:# added by pdchristian to catch decryption error
                file_content = "Empty due to extraction error."  # added by pdchristian to catch decryption error
        elif file.endswith(".txt") or file.endswith(".md") or file.endswith(".markdown"):
            print("indexing " + file)
            f = open(file,'r',encoding='utf-8',errors='ignore')
            file_content = f.read()
            f.close()
        elif file.endswith(".docx"):
            print("indexing " + file)
            file_content = getTextFromWord(file)
        elif file.endswith(".pptx"):
            print("indexing " + file)
            file_content = getTextFromPPTX(file)
        else:
            continue
        text_splitter = TokenTextSplitter(chunk_size=500, chunk_overlap=50)
        texts = text_splitter.split_text(file_content)
        metadata = []
        for i in range(0,len(texts)):
            metadata.append({"path":file})
        qdrant.add_texts(texts,metadatas=metadata)
        len(texts)
    print(onlyfiles)
    print("Finished indexing!")

if __name__ == "__main__":
    arguments = sys.argv
    if len(arguments)>1:
        main_indexing(arguments[1])
    else:
        print("You need to provide a path to folder with documents to index as command line argument")